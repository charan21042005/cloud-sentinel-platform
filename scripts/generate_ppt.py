import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(md_path, output_path):
    if not os.path.exists(md_path):
        print(f"Error: {md_path} not found.")
        return

    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split slides by Markdown horizontal rule (---)
    raw_slides = content.split('---')
    
    prs = Presentation()
    
    # Optional: Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, slide_content in enumerate(raw_slides):
        slide_content = slide_content.strip()
        if not slide_content:
            continue
            
        # Skip the Marp frontmatter
        if 'marp: true' in slide_content:
            continue

        # Extract Title (starting with #)
        title_match = re.search(r'^#\s+(.+)$', slide_content, re.MULTILINE)
        title = title_match.group(1) if title_match else f"Slide {idx}"
        
        # Remove the title and HTML comments from content
        body = re.sub(r'^#\s+.+$', '', slide_content, flags=re.MULTILINE)
        body = re.sub(r'<!--.*?-->', '', body, flags=re.DOTALL)
        body = body.strip()

        # Choose slide layout
        if idx == 1 or "Cloud Sentinel" in title:
            slide_layout = prs.slide_layouts[0] # Title slide
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            title_shape.text = title
            subtitle.text = body[:200] + "..." if len(body) > 200 else body
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = title
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            # Simple parsing for bullet points vs normal text
            lines = body.split('\n')
            first_p = True
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Remove bold asterisks for clean text
                clean_line = line.replace('**', '')
                
                if first_p:
                    p = tf.paragraphs[0]
                    first_p = False
                else:
                    p = tf.add_paragraph()
                
                if clean_line.startswith('- '):
                    p.text = clean_line[2:]
                    p.level = 1
                elif clean_line.startswith('* '):
                    p.text = clean_line[2:]
                    p.level = 1
                else:
                    p.text = clean_line
                    p.level = 0
                    
                p.font.size = Pt(18)

    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}!")

if __name__ == "__main__":
    md_file = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation', 'Cloud_Sentinel_Architecture_Presentation.md')
    out_file = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation', 'Cloud_Sentinel_Architecture_Presentation.pptx')
    create_presentation(md_file, out_file)
